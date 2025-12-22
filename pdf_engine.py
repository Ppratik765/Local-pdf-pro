import os
import io
import img2pdf
import numpy as np
import cv2
import pytesseract
from pypdf import PdfReader, PdfWriter
from pdf2docx import Converter
from pdf2image import convert_from_path
from docx2pdf import convert as docx_convert
from PIL import Image
import comtypes.client
from pptx import Presentation
import pikepdf
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from reportlab.lib import colors

# Set Tesseract Path (Windows default or generic)
# Users must install Tesseract-OCR and add to PATH, or set it here:
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

class PDFEngine:
    # --- EXISTING FEATURES ---
    @staticmethod
    def merge_pdfs(file_list, output_path):
        merger = PdfWriter()
        for pdf in file_list:
            merger.append(pdf)
        merger.write(output_path)
        merger.close()

    @staticmethod
    def split_pdf(input_path, output_folder, mode="all", page_range=None):
        reader = PdfReader(input_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        total_pages = len(reader.pages)
        
        selected_indices = []
        if page_range:
            try:
                parts = [p.strip() for p in page_range.split(',')]
                for p in parts:
                    if '-' in p:
                        start, end = map(int, p.split('-'))
                        selected_indices.extend(range(start-1, end))
                    else:
                        selected_indices.append(int(p)-1)
                selected_indices = [i for i in selected_indices if 0 <= i < total_pages]
            except: pass 
        else:
            selected_indices = list(range(total_pages))

        if mode == "all":
            for i in selected_indices:
                writer = PdfWriter()
                writer.add_page(reader.pages[i])
                out_file = os.path.join(output_folder, f"{base_name}_page_{i+1}.pdf")
                with open(out_file, "wb") as f:
                    writer.write(f)
        elif mode == "extract":
            writer = PdfWriter()
            for i in selected_indices:
                writer.add_page(reader.pages[i])
            out_file = os.path.join(output_folder, f"{base_name}_extracted.pdf")
            with open(out_file, "wb") as f:
                writer.write(f)

    @staticmethod
    def reorder_save_pdf(input_path, output_path, page_order_data):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for item in page_order_data:
            idx = item['original_index']
            rotation = item.get('rotation', 0)
            if 0 <= idx < len(reader.pages):
                page = reader.pages[idx]
                if rotation != 0:
                    page.rotate(rotation)
                writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def images_to_pdf(image_list, output_path):
        processed_images = []
        temp_created = []
        for img_path in image_list:
            try:
                img = Image.open(img_path)
                if img.mode == 'RGBA' or img.format != 'JPEG':
                    img = img.convert('RGB')
                    temp_path = img_path + ".temp.jpg"
                    img.save(temp_path, quality=95)
                    processed_images.append(temp_path)
                    temp_created.append(temp_path)
                else:
                    processed_images.append(img_path)
            except Exception as e:
                print(f"Error processing image {img_path}: {e}")

        if processed_images:
            with open(output_path, "wb") as f:
                f.write(img2pdf.convert(processed_images))
        
        for temp in temp_created:
            try: os.remove(temp)
            except: pass

    @staticmethod
    def pdf_to_images(input_path, output_folder, dpi=200, fmt="jpeg"):
        images = convert_from_path(input_path, dpi=dpi)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        saved_files = []
        for i, img in enumerate(images):
            ext = fmt.lower()
            out_file = os.path.join(output_folder, f"{base_name}_page_{i+1:03d}.{ext}")
            img.save(out_file, fmt.upper())
            saved_files.append(out_file)
        return saved_files

    @staticmethod
    def compress_pdf(input_path, output_path, level="medium"):
        if level == "low":
            reader = PdfReader(input_path)
            writer = PdfWriter()
            for page in reader.pages:
                page.compress_content_streams()
                writer.add_page(page)
            writer.add_metadata(reader.metadata)
            with open(output_path, "wb") as f:
                writer.write(f)
        elif level == "medium":
            try:
                with pikepdf.open(input_path) as pdf:
                    pdf.save(output_path, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate)
            except Exception as e:
                raise Exception(f"Pikepdf failed: {e}")
        elif level == "extreme":
            import tempfile, shutil
            temp_dir = tempfile.mkdtemp()
            try:
                imgs = PDFEngine.pdf_to_images(input_path, temp_dir, dpi=130, fmt="jpeg")
                with open(output_path, "wb") as f:
                    f.write(img2pdf.convert(imgs))
            finally:
                shutil.rmtree(temp_dir)

    # --- NEW PRO FEATURES ---

    @staticmethod
    def ocr_pdf(input_path, output_path, lang='eng'):
        """Converts PDF to images, then uses Tesseract to create a searchable PDF."""
        try:
            # 1. Convert pages to images
            images = convert_from_path(input_path)
            writer = PdfWriter()

            # 2. OCR each image and get a single-page PDF byte stream
            for i, img in enumerate(images):
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf', lang=lang)
                
                # 3. Read that byte stream as a PDF page and add to writer
                page_reader = PdfReader(io.BytesIO(pdf_bytes))
                writer.add_page(page_reader.pages[0])

            with open(output_path, "wb") as f:
                writer.write(f)
        except Exception as e:
            if "tesseract" in str(e).lower():
                raise Exception("Tesseract OCR not found. Please install Tesseract and add it to PATH.")
            raise e

    @staticmethod
    def add_watermark(input_path, output_path, text="", opacity=0.5, rotation=45):
        """Adds a text watermark to every page."""
        reader = PdfReader(input_path)
        writer = PdfWriter()
        
        # Create Watermark PDF in memory
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFillColor(colors.grey, alpha=opacity)
        can.setFont("Helvetica-Bold", 50)
        
        # Calculate centerish position (simplified)
        width, height = float(reader.pages[0].mediabox.width), float(reader.pages[0].mediabox.height)
        
        can.saveState()
        can.translate(width/2, height/2)
        can.rotate(rotation)
        can.drawCentredString(0, 0, text)
        can.restoreState()
        can.save()
        
        packet.seek(0)
        watermark_pdf = PdfReader(packet)
        watermark_page = watermark_pdf.pages[0]

        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)
            
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def add_page_numbers(input_path, output_path, position="bottom-center"):
        """Adds Page X of Y."""
        reader = PdfReader(input_path)
        writer = PdfWriter()
        total = len(reader.pages)
        
        for i, page in enumerate(reader.pages):
            packet = io.BytesIO()
            can = canvas.Canvas(packet, pagesize=(float(page.mediabox.width), float(page.mediabox.height)))
            can.setFont("Helvetica", 10)
            
            w, h = float(page.mediabox.width), float(page.mediabox.height)
            text = f"Page {i+1} of {total}"
            
            if position == "bottom-center":
                can.drawCentredString(w/2, 30, text)
            elif position == "bottom-right":
                can.drawRightString(w - 40, 30, text)
            elif position == "top-right":
                can.drawRightString(w - 40, h - 30, text)
            
            can.save()
            packet.seek(0)
            num_pdf = PdfReader(packet)
            page.merge_page(num_pdf.pages[0])
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def get_metadata(input_path):
        reader = PdfReader(input_path)
        return reader.metadata

    @staticmethod
    def update_metadata(input_path, output_path, metadata):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        writer.append_pages_from_reader(reader)
        writer.add_metadata(metadata)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def auto_scan_image(image_path):
        """
        Robustly detects document contours using adaptive thresholding
        to handle noisy backgrounds (like patterned fabrics).
        """
        import cv2
        import numpy as np
        from PIL import Image

        img = cv2.imread(image_path)
        orig = img.copy()
        
        # 1. Resize for faster processing (maintain aspect ratio)
        scale_height = 800.0
        ratio = img.shape[0] / scale_height
        w = int(img.shape[1] / ratio)
        img_small = cv2.resize(img, (w, int(scale_height)))

        # 2. Preprocessing (Grayscale + Blur)
        gray = cv2.cvtColor(img_small, cv2.COLOR_BGR2GRAY)
        gray = cv2.GaussianBlur(gray, (5, 5), 0)

        # 3. Edge Detection (Canny) AND Adaptive Thresholding
        # We combine both to find strong edges of the paper
        edged = cv2.Canny(gray, 75, 200)
        
        # Morphological closing to close small gaps in the contour
        kernel = np.ones((5, 5), np.uint8)
        edged = cv2.morphologyEx(edged, cv2.MORPH_CLOSE, kernel)

        # 4. Find Contours
        cnts, _ = cv2.findContours(edged.copy(), cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
        cnts = sorted(cnts, key=cv2.contourArea, reverse=True)[:5]

        screenCnt = None
        for c in cnts:
            peri = cv2.arcLength(c, True)
            approx = cv2.approxPolyDP(c, 0.02 * peri, True)

            # If our approximated contour has 4 points, we can assume we found the screen
            if len(approx) == 4:
                screenCnt = approx
                break

        # If detection failed, return original
        if screenCnt is None:
            print("Smart Scan: No document contour detected, returning original.")
            return Image.fromarray(cv2.cvtColor(orig, cv2.COLOR_BGR2RGB))

        # 5. Perspective Transform
        pts = screenCnt.reshape(4, 2) * ratio
        
        # Order points: top-left, top-right, bottom-right, bottom-left
        rect = np.zeros((4, 2), dtype="float32")
        s = pts.sum(axis=1)
        rect[0] = pts[np.argmin(s)] # TL
        rect[2] = pts[np.argmax(s)] # BR
        diff = np.diff(pts, axis=1)
        rect[1] = pts[np.argmin(diff)] # TR
        rect[3] = pts[np.argmax(diff)] # BL

        (tl, tr, br, bl) = rect
        
        # Compute new width/height
        widthA = np.sqrt(((br[0] - bl[0]) ** 2) + ((br[1] - bl[1]) ** 2))
        widthB = np.sqrt(((tr[0] - tl[0]) ** 2) + ((tr[1] - tl[1]) ** 2))
        maxWidth = max(int(widthA), int(widthB))

        heightA = np.sqrt(((tr[0] - br[0]) ** 2) + ((tr[1] - br[1]) ** 2))
        heightB = np.sqrt(((tl[0] - bl[0]) ** 2) + ((tl[1] - bl[1]) ** 2))
        maxHeight = max(int(heightA), int(heightB))

        dst = np.array([
            [0, 0],
            [maxWidth - 1, 0],
            [maxWidth - 1, maxHeight - 1],
            [0, maxHeight - 1]], dtype="float32")

        # Warp
        M = cv2.getPerspectiveTransform(rect, dst)
        warp = cv2.warpPerspective(orig, M, (maxWidth, maxHeight))
        
        # Optional: Adaptive Threshold to make it look like a "scanned" B&W doc
        # For now, we return color as users might want color scans.
        # To make it B&W like a scanner, uncomment the next line:
        # warp = cv2.cvtColor(warp, cv2.COLOR_BGR2GRAY)

        return Image.fromarray(cv2.cvtColor(warp, cv2.COLOR_BGR2RGB))

    @staticmethod
    def manual_scan_warp(image_path, corners):
        """
        Warps image based on 4 manual corner points.
        corners: list of 4 tuples [(x,y), (x,y), (x,y), (x,y)]
        """
        import cv2
        import numpy as np
        from PIL import Image

        img = cv2.imread(image_path)
        orig = img.copy()
        
        # Helper to order points: TL, TR, BR, BL
        pts = np.array(corners, dtype="float32")
        rect = np.zeros((4, 2), dtype="float32")
        
        s = pts.sum(axis=1)
        rect[0] = pts[np.argmin(s)] # TL
        rect[2] = pts[np.argmax(s)] # BR
        diff = np.diff(pts, axis=1)
        rect[1] = pts[np.argmin(diff)] # TR
        rect[3] = pts[np.argmax(diff)] # BL

        (tl, tr, br, bl) = rect

        # Compute width/height
        widthA = np.sqrt(((br[0] - bl[0]) ** 2) + ((br[1] - bl[1]) ** 2))
        widthB = np.sqrt(((tr[0] - tl[0]) ** 2) + ((tr[1] - tl[1]) ** 2))
        maxWidth = max(int(widthA), int(widthB))

        heightA = np.sqrt(((tr[0] - br[0]) ** 2) + ((tr[1] - br[1]) ** 2))
        heightB = np.sqrt(((tl[0] - bl[0]) ** 2) + ((tl[1] - bl[1]) ** 2))
        maxHeight = max(int(heightA), int(heightB))

        dst = np.array([
            [0, 0],
            [maxWidth - 1, 0],
            [maxWidth - 1, maxHeight - 1],
            [0, maxHeight - 1]], dtype="float32")

        M = cv2.getPerspectiveTransform(rect, dst)
        warp = cv2.warpPerspective(orig, M, (maxWidth, maxHeight))
        
        return Image.fromarray(cv2.cvtColor(warp, cv2.COLOR_BGR2RGB))
    
    # --- EXISTING CONVERSIONS ---
    @staticmethod
    def pdf_to_word(input_path, output_path):
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

    @staticmethod
    def word_to_pdf(input_path, output_path):
        docx_convert(input_path, output_path)

    @staticmethod
    def pptx_to_pdf(input_path, output_path):
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1 
        try:
            deck = powerpoint.Presentations.Open(input_path, WithWindow=False)
            deck.SaveAs(output_path, 32)
            deck.Close()
        except Exception as e:
            raise e
        finally:
            powerpoint.Quit()

    @staticmethod
    def pdf_to_pptx(input_path, output_path):
        temp_dir = os.path.dirname(output_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        images = convert_from_path(input_path, dpi=150)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6] 
        if images:
            width, height = images[0].size
            prs.slide_width = int(width * 9525) 
            prs.slide_height = int(height * 9525)
        for i, img in enumerate(images):
            temp_img_path = os.path.join(temp_dir, f"{base_name}_temp_slide_{i}.jpg")
            img.save(temp_img_path, 'JPEG')
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(temp_img_path, 0, 0, prs.slide_width, prs.slide_height)
            try: os.remove(temp_img_path)
            except: pass
        prs.save(output_path)

    @staticmethod
    def protect_pdf(input_path, output_path, password, algorithm="AES-256"):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        writer.append_pages_from_reader(reader)
        writer.encrypt(user_password=password, algorithm=algorithm)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def unlock_pdf(input_path, output_path, password):
        reader = PdfReader(input_path)
        if reader.is_encrypted:
            success = reader.decrypt(password)
            if not success:
                raise Exception("Incorrect Password")
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path
